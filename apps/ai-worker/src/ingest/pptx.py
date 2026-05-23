"""PPTX parser using python-pptx.

One ``Block`` per slide. Modality is ``slide`` so the chunker dispatcher
sends it to ``OneChunkPerBlockStrategy`` — each slide becomes exactly one
chunk and citations point at the right slide index.

Speaker notes are appended to the slide text, separated by a marker line,
so questions about presenter commentary still surface the slide.
"""

from __future__ import annotations

import io
import logging

from pptx import Presentation

from ..agents.contracts import Block, ChunkModality

log = logging.getLogger(__name__)


def parse_pptx(pptx_bytes: bytes) -> list[Block]:
    if not pptx_bytes:
        return []

    blocks: list[Block] = []
    char_cursor = 0
    prs = Presentation(io.BytesIO(pptx_bytes))
    for slide_index, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs).strip()
                    if line:
                        parts.append(line)
        notes = _slide_notes(slide)
        if notes:
            parts.append("\n[notes] " + notes)
        text = "\n".join(parts).strip()
        if not text:
            continue
        blocks.append(
            Block(
                modality=ChunkModality.slide,
                text=text,
                slide=slide_index,
                char_start=char_cursor,
                char_end=char_cursor + len(text),
                meta={"slide_layout": slide.slide_layout.name if slide.slide_layout else None},
            )
        )
        char_cursor += len(text)
    return blocks


def _slide_notes(slide: object) -> str:
    notes_slide = getattr(slide, "notes_slide", None)
    if notes_slide is None:
        return ""
    notes_tf = getattr(notes_slide, "notes_text_frame", None)
    if notes_tf is None:
        return ""
    text = (notes_tf.text or "").strip()
    return text
